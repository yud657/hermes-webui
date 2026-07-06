from __future__ import annotations

import io
import json
import shutil
import subprocess
import zipfile
from pathlib import Path

import pytest

# The Office parsers are optional deps (commented-optional in requirements.txt;
# installed by CI and requirements-dev.txt). On a lean install they are absent,
# so importorskip these modules to skip this file cleanly instead of aborting
# collection for the WHOLE suite with a module-level ImportError (repo idiom —
# see the ~12 other test files that importorskip an optional dependency).
pytest.importorskip("docx")
pytest.importorskip("openpyxl")
pytest.importorskip("pptx")

from docx import Document as DocxDocument
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

import api.office_documents as office_documents
from api.office_documents import (
    CLAIMED_OFFICE_EXTENSIONS,
    is_claimed_office_path,
    preview_office_document,
    save_office_document,
)


ROOT = Path(__file__).resolve().parents[1]
WORKSPACE_JS = (ROOT / "static" / "workspace.js").read_text(encoding="utf-8")
NODE = shutil.which("node")


def _simple_docx_bytes(*paragraphs: str) -> bytes:
    document = DocxDocument()
    for paragraph in paragraphs or ("",):
        document.add_paragraph(paragraph)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _rich_docx_bytes() -> bytes:
    document = DocxDocument()
    document.add_paragraph("Lead paragraph")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "left"
    table.cell(0, 1).text = "right"
    document.add_paragraph("Tail paragraph")
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _formatted_docx_bytes() -> bytes:
    document = DocxDocument()
    run = document.add_paragraph().add_run("Styled text")
    run.bold = True
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _paragraph_style_docx_bytes(style: str) -> bytes:
    document = DocxDocument()
    paragraph = document.add_paragraph("Styled paragraph")
    properties = paragraph._p.get_or_add_pPr()
    paragraph_style = OxmlElement("w:pStyle")
    paragraph_style.set(qn("w:val"), style)
    properties.append(paragraph_style)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _header_docx_bytes() -> bytes:
    document = DocxDocument()
    document.add_paragraph("Body text")
    document.sections[0].header.paragraphs[0].text = "Header text"
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _custom_section_docx_bytes() -> bytes:
    document = DocxDocument()
    document.add_paragraph("Body text")
    document.sections[0].left_margin = Inches(2)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _simple_xlsx_bytes() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Summary"
    sheet["A1"] = "alpha"
    sheet["B1"] = "beta"
    sheet["A2"] = "gamma"
    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def _simple_pptx_bytes() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text = "Office preview"
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()

def _office_zip_bytes(*members: tuple[str, str | bytes]) -> bytes:
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", "<Types/>")
        for name, payload in members:
            archive.writestr(name, payload)
    return buffer.getvalue()


def _multi_slide_pptx_bytes() -> bytes:
    presentation = Presentation()
    for index in range(2):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        box.text = f"Slide {index + 1}"
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def test_office_registry_claims_docx_xlsx_pptx():
    assert CLAIMED_OFFICE_EXTENSIONS == {".docx", ".xlsx", ".pptx"}
    assert is_claimed_office_path("report.docx")
    assert is_claimed_office_path("budget.xlsx")
    assert is_claimed_office_path("deck.pptx")

    docx_preview = preview_office_document("report.docx", _simple_docx_bytes("one", "two"))
    xlsx_preview = preview_office_document("budget.xlsx", _simple_xlsx_bytes())
    pptx_preview = preview_office_document("deck.pptx", _simple_pptx_bytes())

    assert docx_preview["preview_kind"] == "office"
    assert xlsx_preview["preview_kind"] == "office"
    assert pptx_preview["preview_kind"] == "office"
    assert docx_preview["office_format"] == "docx"
    assert xlsx_preview["office_format"] == "xlsx"
    assert pptx_preview["office_format"] == "pptx"
    assert docx_preview["render_mode"] == "code"
    assert xlsx_preview["render_mode"] == "code"
    assert pptx_preview["render_mode"] == "code"
    assert docx_preview["editable"] is True
    assert xlsx_preview["editable"] is False
    assert pptx_preview["editable"] is False


def test_docx_paragraph_projection_round_trips_simple_documents():
    original_bytes = _simple_docx_bytes("alpha", "beta")
    preview = preview_office_document("story.docx", original_bytes)

    assert preview["editable"] is True
    assert preview["content"] == "alpha\nbeta"

    saved_preview, saved_bytes = save_office_document("story.docx", original_bytes, "alpha\nbeta\ngamma")
    round_trip = DocxDocument(io.BytesIO(saved_bytes))

    assert saved_preview["editable"] is True
    assert saved_preview["content"] == "alpha\nbeta\ngamma"
    assert [paragraph.text for paragraph in round_trip.paragraphs] == ["alpha", "beta", "gamma"]


def test_docx_preview_preserves_interleaved_table_order():
    preview = preview_office_document("report.docx", _rich_docx_bytes())

    assert preview["content"] == "Lead paragraph\nTable 1\nleft\tright\nTail paragraph"


def test_docx_with_explicit_normal_style_stays_editable():
    original_bytes = _paragraph_style_docx_bytes("Normal")
    preview = preview_office_document("styled.docx", original_bytes)

    assert preview["editable"] is True
    assert preview["content"] == "Styled paragraph"

    saved_preview, saved_bytes = save_office_document("styled.docx", original_bytes, "Edited paragraph")
    round_trip = DocxDocument(io.BytesIO(saved_bytes))

    assert saved_preview["editable"] is True
    assert saved_preview["content"] == "Edited paragraph"
    assert [paragraph.text for paragraph in round_trip.paragraphs] == ["Edited paragraph"]


def test_docx_with_non_default_paragraph_style_stays_preview_only():
    preview = preview_office_document("heading.docx", _paragraph_style_docx_bytes("Heading1"))

    assert preview["editable"] is False
    assert preview.get("edit_blocked_reason")


def test_docx_preview_truncation_disables_editing(monkeypatch):
    monkeypatch.setattr(office_documents, "MAX_DOCX_PREVIEW_BLOCKS", 1)
    monkeypatch.setattr(
        office_documents,
        "_docx_editability",
        lambda _document: pytest.fail("truncated docx preview should not run full editability scan"),
    )

    preview = preview_office_document("story.docx", _simple_docx_bytes("alpha", "beta"))

    assert preview["truncated"] is True
    assert preview["editable"] is False
    assert preview["edit_blocked_reason"] == "docx preview exceeds safe limits"
    assert office_documents.OFFICE_PREVIEW_TRUNCATED_NOTICE in preview["content"]


def test_docx_preview_char_budget_disables_editing(monkeypatch):
    monkeypatch.setattr(office_documents, "MAX_OFFICE_PREVIEW_CHARS", 5)
    monkeypatch.setattr(
        office_documents,
        "_docx_editability",
        lambda _document: pytest.fail("char-budget truncation should not run full editability scan"),
    )

    preview = preview_office_document("story.docx", _simple_docx_bytes("alphabet", "beta"))

    assert preview["truncated"] is True
    assert preview["editable"] is False
    assert preview["edit_blocked_reason"] == "docx preview exceeds safe limits"
    assert office_documents.OFFICE_PREVIEW_TRUNCATED_NOTICE in preview["content"]

@pytest.mark.parametrize(
    ("path", "member_name", "loader_attr"),
    [
        ("story.docx", "word/document.xml", "_load_docx_document"),
        ("budget.xlsx", "xl/sharedStrings.xml", "_load_workbook_reader"),
        ("deck.pptx", "ppt/media/image1.png", "_load_presentation_ctor"),
    ],
)
def test_office_preview_preflight_rejects_zip_bombs_before_loader(
    monkeypatch, path: str, member_name: str, loader_attr: str
):
    monkeypatch.setattr(office_documents, "MAX_OFFICE_ARCHIVE_TOTAL_UNCOMPRESSED_BYTES", 1024)
    monkeypatch.setattr(
        office_documents,
        loader_attr,
        lambda: lambda *_args, **_kwargs: pytest.fail("office parser should not run after archive preflight rejects the payload"),
    )

    raw = _office_zip_bytes((member_name, "A" * 4096))

    with pytest.raises(ValueError, match="safe archive limits"):
        preview_office_document(path, raw)

def test_office_preview_preflight_rejects_path_traversal_before_loader(monkeypatch):
    monkeypatch.setattr(
        office_documents,
        "_load_docx_document",
        lambda: lambda *_args, **_kwargs: pytest.fail("docx parser should not run after path validation fails"),
    )

    raw = _office_zip_bytes(("../word/document.xml", "payload"))

    with pytest.raises(ValueError, match="safe archive limits"):
        preview_office_document("story.docx", raw)

def test_xlsx_preview_preflight_rejects_oversized_shared_strings_before_loader(monkeypatch):
    monkeypatch.setattr(office_documents, "MAX_XLSX_ARCHIVE_SHARED_STRINGS_BYTES", 64)
    monkeypatch.setattr(
        office_documents,
        "_load_workbook_reader",
        lambda: lambda *_args, **_kwargs: pytest.fail("xlsx parser should not run after sharedStrings preflight rejects the payload"),
    )

    raw = _office_zip_bytes(("xl/sharedStrings.xml", "A" * 256))

    with pytest.raises(ValueError, match="safe archive limits"):
        preview_office_document("budget.xlsx", raw)

def test_xlsx_preview_preflight_measures_actual_bytes_when_zip_claims_zero(monkeypatch):
    class FakeInfo:
        filename = "xl/sharedStrings.xml"
        file_size = 0
        compress_size = 1

        @staticmethod
        def is_dir():
            return False

    class FakeArchive:
        def __enter__(self):
            return self

        def __exit__(self, exc_type, exc, tb):
            return False

        @staticmethod
        def infolist():
            return [FakeInfo()]

        @staticmethod
        def open(_info):
            return io.BytesIO(b"A" * 2048)

    monkeypatch.setattr(office_documents, "MAX_OFFICE_ARCHIVE_TOTAL_UNCOMPRESSED_BYTES", 1024)
    monkeypatch.setattr(office_documents.zipfile, "ZipFile", lambda *_args, **_kwargs: FakeArchive())
    monkeypatch.setattr(
        office_documents,
        "_load_workbook_reader",
        lambda: lambda *_args, **_kwargs: pytest.fail("xlsx parser should not run when actual inflated bytes exceed the preflight cap"),
    )

    with pytest.raises(ValueError, match="safe archive limits"):
        preview_office_document("budget.xlsx", b"placeholder")


def _office_state_block() -> str:
    marker = "if(data.preview_kind==='office'){"
    start = WORKSPACE_JS.find(marker)
    assert start >= 0, "office preview state block not found in static/workspace.js"
    depth = 0
    for idx in range(start, len(WORKSPACE_JS)):
        char = WORKSPACE_JS[idx]
        if char == "{":
            depth += 1
        elif char == "}":
            depth -= 1
            if depth == 0:
                return WORKSPACE_JS[start:idx + 1]
    raise AssertionError("could not find balanced office preview state block")


def _run_office_state_block(payload: dict) -> dict:
    js = (
        "const data = " + json.dumps(payload) + ";\n"
        + "const path = 'report';\n"
        + "let _previewRawContent = null;\n"
        + "let _previewRawContentPath = null;\n"
        + "let _previewServerEditable = null;\n"
        + "let _previewPreviewKind = '';\n"
        + "let _previewOfficeFormat = '';\n"
        + "let _previewSaveRoute = '/api/file/save';\n"
        + _office_state_block()
        + "\nconsole.log(JSON.stringify({_previewRawContent,_previewRawContentPath,_previewServerEditable,_previewPreviewKind,_previewOfficeFormat,_previewSaveRoute}));\n"
    )
    result = subprocess.run([NODE, "-e", js], check=True, capture_output=True, text=True, timeout=30)
    return json.loads(result.stdout.strip())


def test_xlsx_stays_preview_only():
    path = "budget.xlsx"
    raw = _simple_xlsx_bytes()

    preview = preview_office_document(path, raw)

    assert preview["preview_kind"] == "office"
    assert preview["editable"] is False
    assert preview.get("edit_blocked_reason")

    with pytest.raises(ValueError):
        save_office_document(path, raw, "edited text")


def test_xlsx_preview_is_bounded(monkeypatch):
    monkeypatch.setattr(office_documents, "MAX_XLSX_PREVIEW_ROWS_PER_SHEET", 1)

    preview = preview_office_document("budget.xlsx", _simple_xlsx_bytes())

    assert preview["truncated"] is True
    assert office_documents.OFFICE_PREVIEW_TRUNCATED_NOTICE in preview["content"]


def test_xlsx_preview_stops_when_char_budget_is_exhausted(monkeypatch):
    class FakeSheet:
        title = "Summary"
        max_row = 9_999
        max_column = 9_999

        def iter_rows(self, *, values_only, max_row, max_col):
            assert values_only is True
            assert max_row == office_documents.MAX_XLSX_PREVIEW_ROWS_PER_SHEET
            assert max_col == office_documents.MAX_XLSX_PREVIEW_CELLS_PER_SHEET
            yield ("x" * 200, "y" * 200)
            pytest.fail("xlsx preview should stop after the preview budget is exhausted")

    class FakeWorkbook:
        def __init__(self):
            self.worksheets = [FakeSheet()]

        def close(self):
            return None

    monkeypatch.setattr(office_documents, "MAX_OFFICE_PREVIEW_CHARS", 32)
    monkeypatch.setattr(office_documents, "_preflight_office_archive", lambda *_args, **_kwargs: None)
    monkeypatch.setattr(office_documents, "_load_workbook_reader", lambda: lambda *_args, **_kwargs: FakeWorkbook())

    preview = preview_office_document("budget.xlsx", b"placeholder")

    assert preview["truncated"] is True
    assert office_documents.OFFICE_PREVIEW_TRUNCATED_NOTICE in preview["content"]


def test_pptx_stays_preview_only():
    path = "deck.pptx"
    raw = _simple_pptx_bytes()

    preview = preview_office_document(path, raw)

    assert preview["preview_kind"] == "office"
    assert preview["editable"] is False
    assert preview.get("edit_blocked_reason")

    with pytest.raises(ValueError):
        save_office_document(path, raw, "edited text")


def test_pptx_preview_is_bounded(monkeypatch):
    monkeypatch.setattr(office_documents, "MAX_PPTX_PREVIEW_SLIDES", 1)

    preview = preview_office_document("deck.pptx", _multi_slide_pptx_bytes())

    assert preview["truncated"] is True
    assert office_documents.OFFICE_PREVIEW_TRUNCATED_NOTICE in preview["content"]


def test_rich_docx_stays_preview_only():
    path = "rich.docx"
    raw = _rich_docx_bytes()

    preview = preview_office_document(path, raw)

    assert preview["preview_kind"] == "office"
    assert preview["editable"] is False
    assert preview.get("edit_blocked_reason")

    with pytest.raises(ValueError):
        save_office_document(path, raw, "edited text")


def test_formatted_docx_stays_preview_only():
    path = "formatted.docx"
    raw = _formatted_docx_bytes()

    preview = preview_office_document(path, raw)

    assert preview["preview_kind"] == "office"
    assert preview["editable"] is False
    assert preview.get("edit_blocked_reason")

    with pytest.raises(ValueError):
        save_office_document(path, raw, "edited text")

def test_header_docx_stays_preview_only():
    path = "headed.docx"
    raw = _header_docx_bytes()

    preview = preview_office_document(path, raw)

    assert preview["preview_kind"] == "office"
    assert preview["editable"] is False
    assert preview.get("edit_blocked_reason")

    with pytest.raises(ValueError):
        save_office_document(path, raw, "edited text")


def test_custom_section_docx_stays_preview_only():
    path = "custom-section.docx"
    raw = _custom_section_docx_bytes()

    preview = preview_office_document(path, raw)

    assert preview["preview_kind"] == "office"
    assert preview["editable"] is False
    assert preview.get("edit_blocked_reason")

    with pytest.raises(ValueError):
        save_office_document(path, raw, "edited text")


@pytest.mark.skipif(NODE is None, reason="node not on PATH")
def test_workspace_office_state_block_routes_all_office_formats_through_office_save():
    xlsx_state = _run_office_state_block(
        {
            "preview_kind": "office",
            "office_format": "xlsx",
            "editable": False,
            "content": "sheet preview",
        }
    )
    docx_state = _run_office_state_block(
        {
            "preview_kind": "office",
            "office_format": "docx",
            "editable": True,
            "content": "doc preview",
        }
    )

    assert xlsx_state["_previewSaveRoute"] == "/api/file/office-save"
    assert xlsx_state["_previewServerEditable"] is False
    assert xlsx_state["_previewOfficeFormat"] == "xlsx"
    assert docx_state["_previewSaveRoute"] == "/api/file/office-save"
    assert docx_state["_previewServerEditable"] is True


# ---------------------------------------------------------------------------
# Regression coverage for the four blockers found in the v0.51.802 gate review
# (all reproduced against legitimate-but-common inputs the earlier adversarial
# rounds never exercised).
# ---------------------------------------------------------------------------


def _write_only_xlsx_bytes() -> bytes:
    """A workbook produced by openpyxl's streaming writer.

    ``Workbook(write_only=True)`` omits the ``<dimension>`` record, so when the
    preview re-opens it in read_only mode ``sheet.max_row``/``max_column`` are
    ``None``. This is exactly issue #540's own use case (agents streaming large
    sheets) and previously crashed the preview with a ``TypeError`` 500.
    """
    workbook = Workbook(write_only=True)
    sheet = workbook.create_sheet("Streamed")
    sheet.append(["alpha", "beta", "gamma"])
    sheet.append(["one", "two", "three"])
    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def _multi_run_docx_bytes() -> bytes:
    """A docx whose paragraph spans multiple runs with a run-boundary space.

    ``add_run("Hello ") + add_run("world")`` is ubiquitous in real Word files
    (and trivially produced by agent code). Each run is a separate ``<w:t>``
    node; the preview previously ``.strip()``-ed each node and concatenated with
    no separator, corrupting the text to "Helloworld".
    """
    document = DocxDocument()
    paragraph = document.add_paragraph()
    paragraph.add_run("Hello ")
    paragraph.add_run("world")
    document.add_paragraph("Second paragraph.")
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def test_write_only_xlsx_without_dimension_previews_without_crash():
    """BLOCKING 1: dimension-less xlsx (max_row/max_column None) must not 500."""
    raw = _write_only_xlsx_bytes()
    # Must not raise TypeError — the pre-fix code did `min(None, CAP)`.
    preview = preview_office_document("streamed.xlsx", raw)
    assert preview["preview_kind"] == "office"
    assert preview["office_format"] == "xlsx"
    # The streamed cell values must actually appear in the preview.
    assert "alpha" in preview["content"]
    assert "three" in preview["content"]


def test_multi_run_docx_preserves_run_boundary_whitespace():
    """BLOCKING 2 (preview): run-boundary space must survive extraction."""
    raw = _multi_run_docx_bytes()
    preview = preview_office_document("multi-run.docx", raw)
    assert "Hello world" in preview["content"]
    assert "Helloworld" not in preview["content"]


def test_multi_run_docx_round_trips_text_identically_on_save():
    """BLOCKING 2 (save): opening the editor and saving unchanged preview text
    must NOT corrupt the on-disk document. This is the data-loss path: the
    editor textarea is prefilled from the preview text, so a no-op save writes
    the preview text back verbatim."""
    raw = _multi_run_docx_bytes()
    preview = preview_office_document("multi-run.docx", raw)
    assert preview["editable"] is True

    # Simulate a save with the UNMODIFIED preview text (the no-op-save path).
    _saved_preview, saved_bytes = save_office_document(
        "multi-run.docx", raw, preview["content"]
    )

    # Read the saved bytes back and confirm the paragraph text is intact.
    reloaded = DocxDocument(io.BytesIO(saved_bytes))
    paragraph_texts = [p.text for p in reloaded.paragraphs]
    assert "Hello world" in paragraph_texts
    assert "Helloworld" not in paragraph_texts


def test_corrupt_archive_member_raises_value_error_not_unhandled():
    """BLOCKING 3: a valid zip central directory with corrupt member DATA (bad
    CRC / truncated) must raise the module's ValueError (handled as a clean
    error) instead of an unhandled BadZipFile/zlib.error that escapes as a 500."""
    raw = bytearray(_simple_docx_bytes("hello"))
    # Corrupt the deflate stream bytes without touching the central directory,
    # so the archive opens but a member.read() fails a CRC/zlib check. Flip a
    # run of bytes in the middle of the compressed payload region.
    start = len(raw) // 3
    for i in range(start, start + 64):
        raw[i] ^= 0xFF
    with pytest.raises(ValueError):
        preview_office_document("corrupt.docx", bytes(raw))


def test_archive_limit_error_still_raised_over_corrupt_catch():
    """BLOCKING 3 guard: the limit-exceeded ValueError raised INSIDE the member
    read loop must not be swallowed by the new corrupt-member except clause
    (it catches only decompression exceptions, never ValueError)."""
    # A docx whose main document part inflates beyond the per-member cap.
    document = DocxDocument()
    document.add_paragraph("x" * 5_000)
    buffer = io.BytesIO()
    document.save(buffer)
    raw = buffer.getvalue()
    # Temporarily force a tiny member cap so the real (well-formed) archive
    # trips the limit path — proving the limit ValueError propagates.
    original = office_documents.MAX_OFFICE_ARCHIVE_MEMBER_BYTES
    office_documents.MAX_OFFICE_ARCHIVE_MEMBER_BYTES = 10
    try:
        with pytest.raises(ValueError):
            preview_office_document("big.docx", raw)
    finally:
        office_documents.MAX_OFFICE_ARCHIVE_MEMBER_BYTES = original


def _authored_docx_bytes() -> bytes:
    """A docx carrying core-properties metadata and body text (like every
    Word-authored file)."""
    document = DocxDocument()
    document.core_properties.author = "Jane Author"
    document.core_properties.title = "Quarterly Report"
    document.add_paragraph("Body text here")
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def test_unedited_docx_save_preserves_core_properties():
    """Round-trip fidelity: an unedited open->save must NOT wipe author/title.

    The pre-fix save rebuilt the package from python-docx's blank template, so
    metadata/styles were silently lost even when the user changed nothing (the
    editor prefills the textarea from the preview). Rebuilding from the current
    package preserves docProps/styles/theme."""
    raw = _authored_docx_bytes()
    preview = preview_office_document("authored.docx", raw)
    assert preview["editable"] is True

    _saved_preview, saved_bytes = save_office_document(
        "authored.docx", raw, preview["content"]
    )
    reloaded = DocxDocument(io.BytesIO(saved_bytes))
    assert reloaded.core_properties.author == "Jane Author"
    assert reloaded.core_properties.title == "Quarterly Report"
    assert "Body text here" in [p.text for p in reloaded.paragraphs]


def test_oversized_docx_save_content_rejected_before_build():
    """Write-surface DoS guard: content exceeding the editable char/line caps is
    rejected up front, before the quadratic add_paragraph build loop runs."""
    raw = _authored_docx_bytes()
    # Exceed the paragraph-count cap (MAX_DOCX_PREVIEW_BLOCKS).
    too_many_lines = "\n".join("x" for _ in range(office_documents.MAX_DOCX_PREVIEW_BLOCKS + 50))
    with pytest.raises(ValueError):
        save_office_document("authored.docx", raw, too_many_lines)
    # Exceed the char cap with few lines.
    too_many_chars = "a" * (office_documents.MAX_OFFICE_PREVIEW_CHARS + 1)
    with pytest.raises(ValueError):
        save_office_document("authored.docx", raw, too_many_chars)


@pytest.mark.parametrize(
    "paragraphs",
    [
        ["Hello", "", ""],   # trailing blank paragraphs (normal Word convention)
        ["", "World"],        # leading blank paragraph
        ["a", "", "b"],       # interior blank (already worked — guard against regression)
    ],
)
def test_unedited_docx_save_preserves_edge_blank_paragraphs(paragraphs):
    """Round-trip fidelity: leading/trailing blank paragraphs must survive an
    unedited open->save. _finalize_preview_text used to .strip() the whole
    preview, and since the editor textarea is prefilled from the preview text,
    a no-op save silently dropped edge blank paragraphs. The docx preview path
    now preserves edge whitespace (strip_edges=False)."""
    document = DocxDocument()
    for text in paragraphs:
        document.add_paragraph(text)
    buffer = io.BytesIO()
    document.save(buffer)
    raw = buffer.getvalue()

    preview = preview_office_document("edge.docx", raw)
    assert preview["editable"] is True

    _saved_preview, saved_bytes = save_office_document("edge.docx", raw, preview["content"])
    reloaded = [p.text for p in DocxDocument(io.BytesIO(saved_bytes)).paragraphs]
    assert reloaded == paragraphs
